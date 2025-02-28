# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""
import pandas as pd
import numpy as np
import time
import plotly.express as px
import plotly.graph_objects as go
from plotly import tools as tls
from plotly.subplots import make_subplots
import os
from datetime import datetime as dt, timedelta
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
# Import the required tkinter Libraries
import tkinter as tk
from tkinter import *
from tkinter import ttk, filedialog
#from tkinter.filedialog import askopenfile
from tkinter import messagebox
from threading import Thread

# Ignore warnings
import warnings
warnings.simplefilter(action='ignore', category=UserWarning)

#path of the file
def read_files(file_path):
    #List of all files in the folder
    df_list = []
    for file in file_path:
        #Read CSV files
        if file.endswith('.csv'):
            df1 = pd.read_csv(file)
            df_list.append(df1)
        #Read Excel files    
        elif file.endswith('.xlsx'):
            df1 = pd.read_excel(file)
            df_list.append(df1)
        else:
            raise ValueError("Please provide a CSV or Excel file")
    df = pd.concat(df_list, ignore_index=True)
    return df


def pol_test(df):
    #Sort the DataFrame by column Time
    df.sort_values(by='Time', inplace=True)
 
    pol =[]
    ss = [] 
    bp = []
    rbp = []
    ocv = []
    pol_final = []
    ss_final = []
    bp_final = []
    ocv_final = []
    current_phase = []
    current_group = []

    for idx, row in df.iterrows():
        #Check if the phase changes from the previous row
        if row['Phase'] != current_phase:
            #If the previous phase was 'polarization'
            if current_phase == 'Polarization':
                #next_seq = df.iloc[idx+1]['Phase'] if idx+1< len(df) else None
                #if next_seq == 'Steady State' or next_seq == 'Back Pressure Test':
                if row['Phase'] in ['Steady State', 'Back Pressure Test', 'Reverse BPR']:
                    #Append the previous group to the desired phase list
                    pol.append(current_group)    
                    #Append the previous group to the desired phase list
                    #pol.append(current_group)
            elif current_phase == 'Steady State':
                ss.append(current_group)

            elif current_phase == 'Back Pressure Test':
                bp.append(current_group)
                
            elif current_phase == 'Reverse BPR':
                rbp.append(current_group)
                
                
            current_phase = row['Phase']
            current_group = []
        current_group.append(row)
        
    for idx, row in df.iterrows():
        #Check if the phase changes from the previous row
        if row['Phase'] != current_phase:
            if current_phase == 'OCV Decay':
                ocv.append(current_group)
            current_phase = row['Phase']
            current_group = []
        current_group.append(row)

          
    pol_df = pd.concat([pd.DataFrame(group) for group in pol])
    ss_df =  pd.concat([pd.DataFrame(group) for group in ss])
    bp_df = pd.concat([pd.DataFrame(group) for group in bp])
    rbp_df = pd.concat([pd.DataFrame(group) for group in rbp])
    ocv_df = pd.concat([pd.DataFrame(group) for group in ocv])
    bp_df = bp_df[bp_df['Phase'].eq('Back Pressure Test') & (bp_df['Current SP'] == 840)]
    ss_df = ss_df[ss_df['Phase'].eq('Steady State') & (ss_df['Current SP']==840)]
    rbp_df = rbp_df[rbp_df['Phase'].eq('Reverse BPR') & (rbp_df['Current SP'] == 840)]
    ocv_df = ocv_df[ocv_df['Phase'].eq('OCV Decay') & (ocv_df['Current SP']==10)]
    final_df = pd.concat([pol_df, ss_df, bp_df, rbp_df, ocv_df])
    
    return final_df


def raw_data(df):
    new = df.copy()
    #Filter the dataframe by polarization, steady state and back pressure test
    new_df = new[new['Phase'].isin(['Polarization', 'Steady State', 'Back Pressure Test', 'Reverse BPR', 'OCV Decay'])]
    #Drop the unrequired fields
    final_df = new_df.drop(columns=['HYS 101','HYS 401','OXS 101','LVL 101','S2-C8-15',
    'S2-C8-16','S3-C9-1','S3-C9-2','S3-C9-3','S3-C9-4','S3-C9-5','S3-C9-6','S3-C9-7','S3-C9-8','S3-C9-9','S3-C9-10','S3-C9-11','S3-C9-12',
    'S3-C9-13','S3-C9-14','S3-C9-15','S3-C9-16','S3-C10-1','S3-C10-2','S3-C10-3','S3-C10-4','S3-C10-5','S3-C10-6','S3-C10-7','S3-C10-8','S3-C10-9',
    'S3-C10-10','S3-C10-11','S3-C10-12','S3-C10-13','S3-C10-14','S3-C10-15','S3-C10-16','S3-C11-1','S3-C11-2','S3-C11-3','S3-C11-4','S3-C11-5',
    'S3-C11-6','S3-C11-7','S3-C11-8','S3-C11-9','S3-C11-10','S3-C11-11','S3-C11-12','S3-C11-13','S3-C11-14','S3-C11-15','S3-C11-16','S3-C12-1',
    'S3-C12-2','S3-C12-3','S3-C12-4','S3-C12-5','S3-C12-6','S3-C12-7','S3-C12-8','S3-C12-9','S3-C12-10','S3-C12-11','S3-C12-12','S3-C12-13',
    'S3-C12-14','S3-C12-15','S3-C12-16','S4-C13-1','S4-C13-2','S4-C13-3','S4-C13-4','S4-C13-5','S4-C13-6','S4-C13-7','S4-C13-8','S4-C13-9','S4-C13-10',
    'S4-C13-11','S4-C13-12','S4-C13-13','S4-C13-14','S4-C13-15','S4-C13-16','S4-C14-1','S4-C14-2','S4-C14-3','S4-C14-4','S4-C14-5','S4-C14-6','S4-C14-7',
    'S4-C14-8','S4-C14-9','S4-C14-10','S4-C14-11','S4-C14-12','S4-C14-13','S4-C14-14','S4-C14-15','S4-C14-16','S4-C15-1','S4-C15-2','S4-C15-3',
    'S4-C15-4','S4-C15-5','S4-C15-6','S4-C15-7','S4-C15-8','S4-C15-9','S4-C15-10','S4-C15-11','S4-C15-12','S4-C15-13','S4-C15-14','S4-C15-15',
    'S4-C15-16','S4-C16-1','S4-C16-2','S4-C16-3','S4-C16-4','S4-C16-5','S4-C16-6','S4-C16-7','S4-C16-8','S4-C16-9','S4-C16-10','S4-C16-11',
    'S4-C16-12','S4-C16-13','S4-C16-14','S4-C16-15','S4-C16-16','DCDC-0-1 V','DCDC-0-1 I','DCDC-0-2 V','DCDC-0-2 I','DCDC-0-3 V','DCDC-0-3 I','DCDC-0-4 V',
    'DCDC-0-4 I','DCDC-0-5 V','DCDC-0-5 I','DCDC-0-6 V','DCDC-0-6 I','DCDC-0-7 V','DCDC-0-7 I','DCDC-0-8 V','DCDC-0-8 I','DCDC-0-9 V','DCDC-0-9 I','DCDC-0-10 V',
    'DCDC-0-10 I','DCDC-0-11 V','DCDC-0-11 I','DCDC-1-1 V','DCDC-1-1 I','DCDC-1-2 V','DCDC-1-2 I','DCDC-1-3 V','DCDC-1-3 I','DCDC-1-4 V','DCDC-1-4 I',
    'DCDC-1-5 V','DCDC-1-5 I','DCDC-1-6 V','DCDC-1-6 I','DCDC-1-7 V','DCDC-1-7 I','DCDC-1-8 V','DCDC-1-8 I','DCDC-1-9 V','DCDC-1-9 I','DCDC-1-10 V',
    'DCDC-1-10 I','DCDC-1-11 V','DCDC-1-11 I'])
    #Drop the every 9th dummy cell column 
    final_df = final_df.drop(columns=['S1-C1-9','S1-C2-2','S1-C2-11','S1-C3-4','S1-C3-13', 'S1-C4-6','S1-C4-15', 'S2-C5-8','S2-C6-1',
                          'S2-C6-10', 'S2-C7-3','S2-C7-12','S2-C8-5','S2-C8-14',])
    return final_df


def processing_final_df(df):
    t_step = 10
    filtered_df = df.copy()
    filtered_df.reset_index(inplace=True)
    filtered_df['Duration(min)'] = (filtered_df.index*10)/60

    
    renaming_columns = ['S1-C1-1','S1-C1-2','S1-C1-3','S1-C1-4','S1-C1-5','S1-C1-6','S1-C1-7','S1-C1-8','S1-C1-10','S1-C1-11',
    'S1-C1-12','S1-C1-13','S1-C1-14','S1-C1-15','S1-C1-16','S1-C2-1','S1-C2-3','S1-C2-4','S1-C2-5','S1-C2-6','S1-C2-7','S1-C2-8',
    'S1-C2-9','S1-C2-10','S1-C2-12','S1-C2-13','S1-C2-14','S1-C2-15','S1-C2-16','S1-C3-1','S1-C3-2','S1-C3-3','S1-C3-5','S1-C3-6',
    'S1-C3-7','S1-C3-8','S1-C3-9','S1-C3-10','S1-C3-11','S1-C3-12','S1-C3-14','S1-C3-15','S1-C3-16','S1-C4-1','S1-C4-2','S1-C4-3',
    'S1-C4-4','S1-C4-5','S1-C4-7','S1-C4-8','S1-C4-9','S1-C4-10','S1-C4-11','S1-C4-12','S1-C4-13','S1-C4-14','S1-C4-16','S2-C5-1',
    'S2-C5-2','S2-C5-3','S2-C5-4','S2-C5-5','S2-C5-6','S2-C5-7','S2-C5-9','S2-C5-10','S2-C5-11','S2-C5-12','S2-C5-13','S2-C5-14',
    'S2-C5-15','S2-C5-16','S2-C6-2','S2-C6-3','S2-C6-4','S2-C6-5','S2-C6-6','S2-C6-7','S2-C6-8','S2-C6-9','S2-C6-11','S2-C6-12',
    'S2-C6-13','S2-C6-14','S2-C6-15','S2-C6-16','S2-C7-1','S2-C7-2','S2-C7-4','S2-C7-5','S2-C7-6','S2-C7-7','S2-C7-8','S2-C7-9',
    'S2-C7-10','S2-C7-11','S2-C7-13','S2-C7-14','S2-C7-15','S2-C7-16','S2-C8-1','S2-C8-2','S2-C8-3','S2-C8-4','S2-C8-6','S2-C8-7',
    'S2-C8-8','S2-C8-9','S2-C8-10','S2-C8-11','S2-C8-12','S2-C8-13']
    #Drop unneccessry columns
    filtered_df = filtered_df.dropna(subset=['Phase'])
    #Generate new column names sequentially
    new_col_names =[f"Cell {i+1}" for i in range(len(renaming_columns))]
    #Create a dictionary mapping old column
    rename_dict = dict(zip(renaming_columns, new_col_names))
    #Rename the old columns name with new columns names
    filtered_df.rename(columns=rename_dict, inplace=True)
    #Create a list of columns start with 'Cell'
    cell_columns = [col for col in filtered_df.columns if col.startswith('Cell')]
    #Convert the columns to numeric format
    for col in cell_columns:
         filtered_df[col] = pd.to_numeric(filtered_df[col], errors='coerce')
    #Create new columns like min cell voltage, max cell voltage and avg cell voltage performing the operation rowise
    filtered_df['min_CV'] = filtered_df[cell_columns].min(axis=1).round(3)
    filtered_df['max_CV'] = filtered_df[cell_columns].max(axis=1).round(3)
    filtered_df['avg_CV'] = filtered_df[cell_columns].mean(axis=1).round(3)

    return filtered_df

#%% this function makes a filename from a fig title
def get_fn_from_fig_title(fig):
    file_nm = fig['layout']['title']['text'].replace(" ", "_")
    file_nm = file_nm.replace("_-_", "-")
    file_nm = file_nm.replace("_&_", "_")
    file_nm = file_nm.lower()
    return file_nm

def plotting_overall_data(df):
    # make sure the current is greater than 0.0049
    df = df.loc[df['Stack Current'] != 0, :]
    filtered_df = df.copy()

    fig= go.Figure()

    fig.add_trace(go.Scatter(x=filtered_df['Duration(min)'], y=filtered_df['Stack Current'], mode='lines', name='Current', line=dict(color='black')))

    fig.add_trace(go.Scatter(x=filtered_df['Duration(min)'], y=filtered_df['PRT 401'], mode='lines', name='Pressure[bar]',yaxis='y2',line=dict(color='orange')))
              
    fig.add_trace(go.Scatter(x=filtered_df['Duration(min)'], y=filtered_df['min_CV'], mode='lines', name='Min CV', yaxis='y3', line=dict(color='green')))
    fig.add_trace(go.Scatter(x=filtered_df['Duration(min)'], y=filtered_df['max_CV'], mode='lines', name='Max CV', yaxis='y3', line=dict(color='blue')))
    fig.add_trace(go.Scatter(x=filtered_df['Duration(min)'], y=filtered_df['avg_CV'], mode='lines', name='Avg CV', yaxis='y3', line=dict(color='deeppink')))

    #Iterate through each phase and plot the data
    for phase in filtered_df['Phase'].unique():
        phase_df = filtered_df[filtered_df['Phase']==phase]
        fig.add_vline(x=phase_df['Duration(min)'].iloc[-1], line_width=5, line_dash='solid', line_color='black'), #annotation_text=f"{phase}")
     
        phase_title = f"{phase}"#" ({int(phase_df['Duration(min)'].iloc[0])}-{int(phase_df['Duration(min)'].iloc[-1])}min)"
        phase_title_x = phase_df['Duration(min)'].iloc[len(phase_df)//2]
        fig.add_annotation(text=phase_title, x=phase_title_x, y=1.02, xref='x', yref='paper', showarrow=False, font=dict(size=12))
    

    fig.update_layout(
        title='All Phases',
        xaxis = dict(title='Duration [min]', domain=[0, 0.90]),
        yaxis=dict(title='Current[A]',  showgrid=False),
        yaxis2=dict(title='Pressure[bar]',overlaying='y', side='right', anchor="x"),
        yaxis3=dict(title='Voltage[V]', overlaying='y', side='right', anchor="free", position=1.0, range=[1.4, 2.2]),
        legend=dict(x=0, y =1.1, bgcolor = 'rgba(0,0,0,0)', orientation='h'),
        height=600, width=1000,
        hovermode='closest',
        template = 'simple_white',
        font = dict(size = 18)
    )   
    
    fig.update_layout(
        yaxis = dict(range=[100, 1200], tickmode='array', tickvals=list(range(100, 1300, 100))))
         
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))

    return fig.show()

def plotting_polarization_data(df):
    filtered_df = df.copy()
    #Filter the original DataFrame where phase is polarization
    polarization_df = filtered_df[filtered_df['Phase']=='Polarization']
    polarization_df = filtered_df[['Current SP', 'Stack Current','min_CV','max_CV','avg_CV']]
    polarization_df = polarization_df[polarization_df.groupby('Current SP')['Current SP'].transform('size')>2]
    # make sure the current is greater than 0.0049
    polarization_df = polarization_df.loc[polarization_df['Stack Current'] > 0, :]
    #polarization_df = polarization_df.groupby(['Current SP']).apply(lambda x: x.iloc[1:-1])
    polarization_df.reset_index(drop=True, inplace=True)
    polarization_df['Current SP'] = pd.to_numeric(polarization_df['Current SP'], errors='coerce')
    polarization_df = polarization_df.groupby(['Current SP']).mean().reset_index()
    polarization_df['j/cm2'] = polarization_df['Stack Current'].div(702.25).round(3)
    polarization_df = polarization_df[polarization_df['Current SP']!= 840]



    fig=go.Figure()
             
    fig.add_trace(go.Scatter(x=polarization_df['j/cm2'], y=polarization_df['min_CV'], mode='lines+markers', name='Min CV', line=dict(color='green')))
    fig.add_trace(go.Scatter(x=polarization_df['j/cm2'], y=polarization_df['max_CV'], mode='lines+markers', name='Max CV', line=dict(color='orange')))
    fig.add_trace(go.Scatter(x=polarization_df['j/cm2'], y=polarization_df['avg_CV'], mode='lines+markers', name='Avg CV', line=dict(color='purple')))

    fig.update_layout(
        title = 'Polarization Chart',
        xaxis = dict(title='Current Density[A/cm2]', rangemode='tozero'),
        yaxis=dict(title='Voltage[V]',showgrid=False),
        width=1000, height=600,
        legend=dict(x=0, y =1.1, bgcolor = 'rgba(0,0,0,0)', orientation='h'),
        template = 'simple_white',
        font = dict(size = 18)
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))

    return fig.show()

def plotting_BPR_plot(df):
    filtered_df = df.copy()
    back_pressure_df =filtered_df[filtered_df['Phase'].isin(['Back Pressure Test', 'Reverse BPR'])]
    back_pressure_df = back_pressure_df.loc[back_pressure_df['Stack Current'] > 0, :]
    back_pressure_df = back_pressure_df.loc[back_pressure_df['min_CV'] > 0, :]
    back_pressure_df = back_pressure_df.loc[back_pressure_df['max_CV'] > 0, :]
    back_pressure_df = back_pressure_df.loc[back_pressure_df['avg_CV'] > 0, :]
    t_step=10
    back_pressure_df.reset_index(inplace=True)
    back_pressure_df.drop(columns='Duration(min)', inplace=True)
    back_pressure_df['Duration(min)'] = (back_pressure_df.index*10)/60
    
    fig=tls.make_subplots(rows=2, cols=1, shared_xaxes=True, vertical_spacing=0.12, horizontal_spacing=0.15)
    fig.add_trace(go.Scatter(x=back_pressure_df['Duration(min)'], y=back_pressure_df['Stack Current'], mode='lines', name='Current',line=dict(color='black')), row=1, col=1)
    fig.add_trace(go.Scatter(x=back_pressure_df['Duration(min)'], y=back_pressure_df['min_CV'], mode='lines',name='Min CV',line=dict(color='green')), row=1, col=1)
    fig.add_trace(go.Scatter(x=back_pressure_df['Duration(min)'], y=back_pressure_df['max_CV'], mode='lines',name='Max CV',line=dict(color='blue')), row=1, col=1)
    fig.add_trace(go.Scatter(x=back_pressure_df['Duration(min)'], y=back_pressure_df['avg_CV'], mode='lines', yaxis='y2',name='Avg CV',line=dict(color='deeppink')), row=1, col=1)
    
    fig.add_trace(go.Scatter(x=back_pressure_df['Duration(min)'], y=back_pressure_df['COS 101'], mode='lines', name='Conductivity',line=dict(color='purple')),   row=2, col=1)
    fig.add_trace(go.Scatter(x=back_pressure_df['Duration(min)'], y=back_pressure_df['HYS 102'], mode='lines', name='h2sensor',line=dict(color='darkorange')),row=2, col=1)
    fig.add_trace(go.Scatter(x=back_pressure_df['Duration(min)'], y=back_pressure_df['PRT 401'], mode='lines', name='Pressure[bar]',line=dict(color='red')), row=2, col=1)

   
    fig.update_layout(
    title = 'Back Pressure Test',
    yaxis=dict(title='Current[A]',range = [100, 1200]),
    yaxis2= dict(title='COS[mu/cm]', rangemode='tozero'),
    xaxis2 = dict(title='Duration [min]'),
    legend=dict(x=0, y =1.1, bgcolor = 'rgba(0,0,0,0)', orientation='h'),
    template = 'simple_white',
    height=600, width=800,
    font = dict(size = 14)
    )   
    with fig.batch_update():
        fig.data[1].update(yaxis='y3')
        fig.layout.update(yaxis3=dict(overlaying='y', side='right',range=[1.6, 2], anchor='x', title='Voltage[V]'))
        fig.data[2].update(yaxis='y3')
        fig.layout.update(yaxis3=dict(overlaying='y', side='right',range=[1.6, 2], anchor='x', title='Voltage[V]'))
        fig.data[3].update(yaxis='y3')
        fig.layout.update(yaxis3=dict(overlaying='y', side='right',range=[1.6, 2], anchor='x', title='Voltage[V]'))

        fig.data[5].update(yaxis='y4')
        fig.layout.update(yaxis4=dict(overlaying='y2', side='right',rangemode='tozero', anchor='x', title='Pressure[bar] & H2 Sensor[%]'))
        fig.data[6].update(yaxis='y4')
        fig.layout.update(yaxis4=dict(overlaying='y2', side='right',rangemode='tozero', anchor='x'))
        # export the figures
    #fig.to_html('bpr.html')
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

   

def plotting_steady_state_vol_dis(df):
    filtered_df = df.copy()
    box_df= filtered_df[filtered_df['Phase']=='Steady State']
    box_df =box_df[box_df['Current SP']==840]
    no_of_cells = 112
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in box_df.columns if col.startswith('Cell')]
    box_final = box_df[box_cells]
    colors = px.colors.qualitative.Dark24
    box_final_df = box_final.copy()

    bundles= [box_final_df.iloc[:, i:i+8] for i in range(0, 112, 8)]
    bundle_means =[bundle.mean(axis=0) for bundle in bundles]
    print(bundle_means)

    bundle_numbers = np.arange(1, len(bundles)+1)
    box_fig = go.Figure()
    for i, bundle_mean in enumerate(bundle_means):
        box_fig.add_trace(go.Box(x=[i+1]*len(bundle_mean),y=bundle_mean, line=dict(color=colors[i]),name=f'Bundle {i+1}', showlegend=False))
   
    scatter_fig = go.Figure()
    for i, bundle_mean in enumerate(bundle_means):
        marker_sizes = [150*abs(x-bundle_mean.mean()) for x in bundle_mean]
        scatter_fig.add_trace(go.Scatter(x=[i+1]*len(bundle_mean), y=bundle_mean, mode='markers',
                                marker=dict(size=marker_sizes, color='black', line=dict(color='black', width=1)), name=f'Bundle {i+1}',  showlegend=False))

    fig = go.Figure(data=box_fig.data + scatter_fig.data) 
    fig.update_layout(
    title='Steady State (840A) Voltage Distribution',
    xaxis=dict(title='Bundle Numbers'),
    yaxis=dict(title='Voltage[V]'),
    width=1000, height=600,
    template = 'simple_white',
    font = dict(size = 18))
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def steady_state_vol_dis_by_pressure(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    ss_df = filtered_df[filtered_df['Phase'].isin(['Steady State'])]

    box_cells = [col for col in ss_df.columns if col.startswith('Cell')]

    new_ss = ss_df[box_cells]
    ss_means = new_ss[box_cells].mean(axis=0)
    print(ss_means)

    box_fig = go.Figure()
    box_fig.add_trace(go.Box(x=['Steady State']*len(ss_means), y=ss_means,line=dict(color='#FF0000'), name='Steady State', showlegend=False))
    
    scatter_fig = go.Figure()
    marker_sizes = [150*abs(x - ss_means.mean()) for x in ss_means]
    scatter_fig.add_trace(go.Scatter(x=['Steady State']*len(ss_means), y=ss_means, mode='markers',
                         marker=dict(size=marker_sizes, color='black'),  name='Steady State',  showlegend=False))

    bp_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    box_cells1 = [col for col in bp_df.columns if col.startswith('Cell')]
    box_cells1.append('PRT 401')
    new = bp_df[box_cells1]
    colors = ['#FF0000', '#FFA500', '#FFFF00', '#008000', '#0000FF', '#4B0082', '#EE82EE', '#FFC0CB', '#1f77b4', '#808080', '#bcbd22']
    cell_columns = [col for col in bp_df.columns if col.startswith('Cell')]

    bundles = {}
    for bp, group_df in new.groupby('PRT 401'):
        bundles[bp] = group_df[cell_columns].values

    bundle_means = {}
    for bp, bundle in bundles.items():
        bundle_means[bp] = np.mean(bundle, axis=0)
        print(bundle_means)

    for bp, bundle_mean in bundle_means.items():
        box_fig.add_trace(go.Box(x=[bp]*len(bundle_mean), y=bundle_mean, line=dict(color=colors[bp]), name=bp, showlegend=False))
    
    for bp, bundle_mean in bundle_means.items():
        marker_sizes = [150*abs(x - bundle_mean.mean()) for x in bundle_mean]
        scatter_fig.add_trace(go.Scatter(x=[bp]*len(bundle_mean), y=bundle_mean, mode='markers',
                         marker=dict(size=marker_sizes, color='black'),  name=bp,  showlegend=False))

    fig = go.Figure(data=box_fig.data + scatter_fig.data) 
    

    fig.update_layout(
    title='Steady State Voltage Distribution by Pressure (All Cells)',
    xaxis=dict(title='Back Pressures'),
    yaxis=dict(title=' Cell Voltage[V]'),
    width=1000, height=600,
    template = 'simple_white',
    font = dict(size = 18))

        # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def back_pressure_by_bundle10(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    bpr_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    
    no_of_cells = 112
    
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in bpr_df.columns if col.startswith('Cell')]
    box_final = bpr_df[box_cells]
    colors = px.colors.qualitative.Dark24
    box_final_df = box_final.copy()

    bundles= [box_final_df.iloc[:, i:i+8] for i in range(0, 112, 8)]
    bundle_means =[bundle.mean(axis=0) for bundle in bundles]
    print(bundle_means)
    bundle_numbers = np.arange(1, len(bundles)+1)
    box_fig = go.Figure()
    for i, bundle_mean in enumerate(bundle_means):

        box_fig.add_trace(go.Box(x=[i+1]*len(bundle_mean),y=bundle_mean, line=dict(color=colors[i]),name=f'Bundle {i+1}', showlegend=False))
  
    scatter_fig = go.Figure()
    for i, bundle_mean in enumerate(bundle_means):
        marker_sizes = [150*abs(x-bundle_mean.mean()) for x in bundle_mean]
        scatter_fig.add_trace(go.Scatter(x=[i+1]*len(bundle_mean), y=bundle_mean, mode='markers',
                                marker=dict(size=marker_sizes, color='black', line=dict(color='black', width=1)), name=f'Bundle {i+1}',  showlegend=False))

    fig = go.Figure(data=box_fig.data + scatter_fig.data) 
    fig.update_layout(
    title='10bar Voltage Distribution',
    xaxis=dict(title='Bundle Numbers'),
    yaxis=dict(title='Voltage[V]'),
    width=1000, height=600,
    template = 'simple_white',
    font = dict(size = 18))
    fig.update_layout(
    xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    
      
    return fig.show()

def back_pressure_by_bundle13(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    bpr_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([13])]
    
    no_of_cells = 112
    
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in bpr_df.columns if col.startswith('Cell')]
    box_final = bpr_df[box_cells]
    colors = px.colors.qualitative.Dark24
    box_final_df = box_final.copy()

    bundles= [box_final_df.iloc[:, i:i+8] for i in range(0, 112, 8)]
    bundle_means =[bundle.mean(axis=0) for bundle in bundles]
    print(bundle_means)
    bundle_numbers = np.arange(1, len(bundles)+1)
    box_fig = go.Figure()
    for i, bundle_mean in enumerate(bundle_means):

        box_fig.add_trace(go.Box(x=[i+1]*len(bundle_mean),y=bundle_mean, line=dict(color=colors[i]),name=f'Bundle {i+1}', showlegend=False))
  
    scatter_fig = go.Figure()
    for i, bundle_mean in enumerate(bundle_means):
        marker_sizes = [150*abs(x-bundle_mean.mean()) for x in bundle_mean]
        scatter_fig.add_trace(go.Scatter(x=[i+1]*len(bundle_mean), y=bundle_mean, mode='markers',
                                marker=dict(size=marker_sizes, color='black', line=dict(color='black', width=1)), name=f'Bundle {i+1}',  showlegend=False))

    fig = go.Figure(data=box_fig.data + scatter_fig.data) 
    fig.update_layout(
    title='13bar Voltage Distribution',
    xaxis=dict(title='Bundle Numbers'),
    yaxis=dict(title='Voltage[V]'),
    width=1000, height=600,
    template = 'simple_white',
    font = dict(size = 18))
    fig.update_layout(
    xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    
      
    return fig.show()

def plot_OCV_decay(df):
    filtered_df = df.copy()
    fig = go.Figure()

    color_cycle = px.colors.qualitative.Dark24

    new_df = filtered_df[(filtered_df['Phase'] == 'OCV Decay')]
    
    print(new_df[['Duration(min)', 'Phase']])
    all_cells = [f'Cell {i}' for i in range(1, 113)]
    df_f = new_df[all_cells]
    df_final = df_f.loc[:, 'Cell 1':'Cell 112']
    #df_final = df_final[(df_final>0).all(axis=1)]
    #select columns
    # selected_columns = []
    # for col in new_df.columns:
    #     if col.startswith('Cell'):
    #         if new_df[col].nunique() > 1:
    #             selected_columns.append(col)

    for i, col in enumerate(all_cells):
        
        duration = new_df['Duration(min)']
        cell_voltages = df_final[col]
        #Add overall line to the plot without showing in the legend
        fig.add_trace(go.Scatter(x=duration, y=cell_voltages, mode='lines+markers', name=f'{col}' , showlegend=True))
              
        fig.update_layout(
            title='OCV Decay Chart',
            xaxis_title='Duration (min)',
            yaxis_title='Cell Voltage [V]',
            width=1400, height=800,
            template = 'simple_white',
            font=dict(size=16),
            legend=dict(orientation = 'h', yanchor='bottom', y=0.2, xanchor='left', x=1),
            )  
              
     

    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
      
    return fig.show()

def plotting_bundle1_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 0:8]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle1',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle2_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 8:16]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle2',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle3_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 16:24]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle3',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle4_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 24:32]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle4',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle5_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 32:40]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle5',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle6_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 40:48]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle6',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle7_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 48:56]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle7',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle8_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 56:64]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle8',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()


def plotting_bundle9_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 64:72]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle9',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()


def plotting_bundle10_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 72:80]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle10',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle11_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 80:88]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle11',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()


def plotting_bundle12_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 88:96]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle12',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle13_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 96:104]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle13',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()

def plotting_bundle14_vol_dis(df):
    filtered_df = df.copy()
    filtered_df["PRT 401"] = filtered_df["PRT 401"].apply(lambda x: round(float(x)))
    filtered_df = filtered_df[filtered_df['Phase'].eq('Back Pressure Test') & filtered_df['PRT 401'].isin([10])]
    colors = px.colors.qualitative.Dark24
    box_cells = [col for col in filtered_df.columns if col.startswith('Cell')]
    box_df = filtered_df[box_cells].copy()
    colors = px.colors.qualitative.Dark24
    box_final = box_df.iloc[:, 104:112]
    all_cell_mean = box_final.mean(axis=0)
    print(all_cell_mean)
    min_val = box_final.min()
    max_val = box_final.max()
    avg_val = box_final.mean()
    
    fig = go.Figure()
    for i, cell in enumerate(box_final.columns):
        fig.add_trace(go.Box(x=[i+1]*len(box_final), y=box_final[cell], line=dict(color=colors[i]), name=f'{cell} Min:{min_val[cell]:.3f} V, Max:{max_val[cell]:.3f} V, Avg:{avg_val[cell]:.3f} V'))
        
   
   
    for i, cell in enumerate(box_final.columns):
        col_mean = all_cell_mean[i]
        marker_sizes = [500*abs(x-col_mean) for x in box_final[cell]]
        marker_colors = ['green' if 1.65<= val <= 1.85 else 'black' for val in box_final[cell]]
        fig.add_trace(go.Scatter(x=[i+1]*len(box_final), y=box_final[cell], mode='markers',
                                marker=dict(size=marker_sizes, color=marker_colors, line=dict(color=marker_colors, width=1)) , showlegend=False))

    
    fig.update_layout(
    title='10 Bar Voltage Distribution for Bundle14',
    xaxis=dict(title='Cell Numbers'),
    yaxis=dict(title='Voltage[V]', tickfont=dict(size=16)),
    width=1200, height=600,
    template = 'simple_white',
    font = dict(size = 18),
    boxgroupgap=0.1)
    fig.update_layout(
        xaxis = dict(
        tickmode = 'linear',
        tick0 = 1,
        dtick = 1
        )
    )
    # export the figures
    fig.write_image('{}/{}.png'.format(image_folder_name, get_fn_from_fig_title(fig)), engine="orca")
    fig.write_html('{}/{}.html'.format(image_folder_name, get_fn_from_fig_title(fig)), include_plotlyjs=('cdn'))
    return fig.show()


# create an images folder to store all static images
#output_dir = r"D:\PST_Data_Log\Jason stacks\U201"


# Read the file
#file_path =  [r"D:\PST_Data_Log\Jason stacks\U195\DATA_PST_U195-112_2024_03_23.csv",
 #            r"D:\PST_Data_Log\Jason stacks\U195\DATA_PST_U195-112_2024_03_24.csv"]

def create_pptx(file_path, image_folder_name, data, filtered_df, final_df, output_dir):
    # load the template presentation
    prs = Presentation('PST_Report.pptx')
    # now update the slide
    slide = prs.slides[0]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(1.3)
    height = Inches(6.0)
    width = Inches(9.0)
    # place the image
    pic = slide.shapes.add_picture('{}/all_phases.png'.format(image_folder_name), row, column,  width=width, height=height)
    pic.click_action.hyperlink.address = '{}/all_phases.html'.format(image_folder_name)
    # update header table
    tables = [shape for shape in slide.shapes if shape.has_table]
    for table in tables:
        num_rows = len(table.table.rows)
        num_cols = len(table.table.columns)
        if num_rows ==1 and num_cols==2:
            # update the stack id cell
            stack_id_cell = table.table.cell(0,1)
            stack_id_cell.text = str(filtered_df['Stack ID'].iloc[-1])
            for paragraph in stack_id_cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)
        elif num_rows == 1 and num_cols == 6:
            duration_cell = table.table.cell(0,1)
            duration_cell.text = '{:.2f} mins'.format(filtered_df['Duration(min)'].iloc[-1])
            start_date = table.table.cell(0,3)
            start_date.text = str(filtered_df['Time'].iloc[0])
            end_date = table.table.cell(0,5)
            end_date.text = str(filtered_df['Time'].iloc[-1])
            for row in table.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(14)
        elif num_rows == 14 and num_cols == 6:
            for phase in filtered_df['Phase'].unique():
                phase_df = filtered_df[filtered_df['Phase']==phase]
                
                if phase == 'Polarization':
                    #phase_df['Current SP'] = phase_df['Current SP'].astype(int)
                    last_row = phase_df[phase_df['Current SP']== 1000].iloc[-1]
                    max_current = table.table.cell(1,1)
                    max_current.text = f"{str(phase_df['Current SP'].max())} A"
                    cell_columns = last_row.filter(like='Cell')
                    cell_col = cell_columns.astype(float)
                    max_cell_value = cell_col.max()
                    max_cell_vol = table.table.cell(2,1)
                    max_cell_vol.text = f"{str(max_cell_value)} V"
                    max_cell_col = cell_col.idxmax()
                    max_cell_column = table.table.cell(3,1)
                    max_cell_column.text = str(max_cell_col)
                    min_cell_value = cell_col.min()
                    min_cell_vol = table.table.cell(4,1)
                    min_cell_vol.text = f"{str(min_cell_value)} V"
                    min_cell_col = cell_col.idxmin()
                    max_cell_column = table.table.cell(5,1)
                    max_cell_column.text = str(min_cell_col)
                    avg_vol = table.table.cell(6,1)
                    avg_vol.text = f"{str(phase_df[phase_df['Current SP']== 1000]['avg_CV'].iloc[-1])} V"
                    water_temp = table.table.cell(7,1)
                    water_temp.text = f"{str(phase_df[phase_df['Current SP']== 1000]['TTC 102'].iloc[-1])} °C"
                    inlet_h2_pre = table.table.cell(8,1)
                    inlet_h2_pre.text = f"{str(phase_df[phase_df['Current SP']== 1000]['PRT 104'].iloc[-1])} bar"
                    cos = table.table.cell(9,1)
                    cos.text = f"{str(phase_df[phase_df['Current SP']== 1000]['COS 101'].iloc[-1])} µS/cm"
                    inherent_pre = table.table.cell(10,1)
                    inherent_pre.text= f"{str(phase_df[phase_df['Current SP']== 1000]['PRT 401'].iloc[-1])} bar"
                    hys1 = table.table.cell(11,1)
                    hys1.text_frame.text = f"{str(phase_df[phase_df['Current SP']== 1000]['HYS 102'].max())} %"
                    hys2 = table.table.cell(12,1)
                    hys2.text_frame.text = f"{str(phase_df[phase_df['Current SP']== 1000]['HYS 102'][-10:].mean())} %"
                    hys3 = table.table.cell(13,1)
                    hys3.text_frame.text = f"{str(phase_df[phase_df['Current SP']== 1000]['HYS 501'].max())} %"
                    for row in table.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(12)
                elif phase == 'Steady State':
                    num_cols = phase_df.select_dtypes(include=['number'])
                    print(num_cols.columns)
                    last_row = num_cols[num_cols['Current SP']==840].mean().round(3)
  
                    max_current = table.table.cell(1,2)
                    max_current.text = f"{str(phase_df['Current SP'].max())} A"
                    cell_columns = last_row.filter(like='Cell')
                    cell_col = cell_columns.astype(float)
                    max_cell_value = cell_col.max()
                    max_cell_vol = table.table.cell(2,2)
                    max_cell_vol.text = f"{str(max_cell_value)} V"
                    max_cell_col = cell_col.idxmax()
                    max_cell_column = table.table.cell(3,2)
                    max_cell_column.text = str(max_cell_col)
                    min_cell_value = cell_col.min()
                    min_cell_vol = table.table.cell(4,2)
                    min_cell_vol.text = f"{str(min_cell_value)} V"
                    min_cell_col = cell_col.idxmin()
                    max_cell_column = table.table.cell(5,2)
                    max_cell_column.text = str(min_cell_col)
                    avg_vol = table.table.cell(6,2)
                    avg_vol.text = f"{str(phase_df[phase_df['Current SP']==840]['avg_CV'].mean().round(3))} V"
                    water_temp = table.table.cell(7,2)
                    water_temp.text = f"{str(phase_df[phase_df['Current SP']==840]['TTC 102'].mean().round(1))} °C"
                    inlet_h2_pre = table.table.cell(8,2)
                    inlet_h2_pre.text = f"{str(phase_df[phase_df['Current SP']==840]['PRT 104'].astype(float).mean().round(2))} bar"
                    cos = table.table.cell(9,2)
                    cos.text = f"{str(phase_df[phase_df['Current SP']==840]['COS 101'].mean().round(1))} µS/cm"
                    inherent_pre = table.table.cell(10,2)
                    inherent_pre.text= f"{str(phase_df[phase_df['Current SP']==840]['PRT 401'].mean().round(2))} bar"
                    hys1 = table.table.cell(11,2)
                    hys1.text_frame.text = f"{str(phase_df[phase_df['Current SP']== 840]['HYS 102'].max())} %"
                    hys2 = table.table.cell(12,2)
                    hys2.text_frame.text = f"{str(phase_df[phase_df['Current SP']== 840]['HYS 102'][-10:].mean())} %"
                    hys3 = table.table.cell(13,2)
                    hys3.text_frame.text = f"{str(phase_df[phase_df['Current SP']== 840]['HYS 501'].max())} %"
                    for row in table.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(12)
                elif phase == 'Back Pressure Test':
                    phase_df["PRT 401"] = phase_df["PRT 401"].apply(lambda x: round(float(x)))
                    phase_df['Current SP'] = phase_df['Current SP'].astype(int)
                    #print(phase_df.to_csv('bpr217valuecheck.csv'))
                    fil = phase_df[phase_df['PRT 401']==10]
                    if not fil.empty:
                        num_cols = phase_df.select_dtypes(include=['number'])
                        last_row = num_cols[num_cols['PRT 401']==10].mean().round(3)
                        max_current = table.table.cell(1,3)
                        max_current.text = f"{str(phase_df['Current SP'].max())} A"
                        cell_columns = last_row.filter(like='Cell')
                        cell_col = cell_columns.astype(float)
                        max_cell_value = cell_col.max()
                        max_cell_vol = table.table.cell(2,3)
                        max_cell_vol.text = f"{str(max_cell_value)} V"
                        max_cell_col = cell_col.idxmax()
                        max_cell_column = table.table.cell(3,3)
                        max_cell_column.text = str(max_cell_col)
                        min_cell_value = cell_col.min()
                        min_cell_vol = table.table.cell(4,3)
                        min_cell_vol.text = f"{str(min_cell_value)} V"
                        min_cell_col = cell_col.idxmin()
                        max_cell_column = table.table.cell(5,3)
                        max_cell_column.text = str(min_cell_col)
                        avg_vol = table.table.cell(6,3)
                        avg_vol.text = f"{str(phase_df[phase_df['PRT 401']==10]['avg_CV'].mean().round(3))} V"
                        water_temp = table.table.cell(7,3)
                        water_temp.text = f"{str(phase_df[phase_df['PRT 401']==10]['TTC 102'].mean().round(1))} °C"
                        inlet_h2_pre = table.table.cell(8,3)
                        inlet_h2_pre.text = f"{str(phase_df[phase_df['PRT 401']==10]['PRT 104'].astype(float).mean().round(2))} bar"
                        cos = table.table.cell(9,3)
                        cos.text = f"{str(phase_df[phase_df['PRT 401']==10]['COS 101'].mean().round(1))} µS/cm"
                        inherent_pre = table.table.cell(10,3)
                        inherent_pre.text= f"{str(phase_df[phase_df['PRT 401']==10]['PRT 401'].mean().round(2))} bar"
                        hys1 = table.table.cell(11,3)
                        hys1.text_frame.text = f"{str(phase_df[phase_df['PRT 401']==10]['HYS 102'].max())} %"
                        hys2 = table.table.cell(12,3)
                        hys2.text_frame.text = f"{str(phase_df[phase_df['PRT 401']==10]['HYS 102'][-10:].mean())} %"
                        hys3 = table.table.cell(13,3)
                        hys3.text_frame.text = f"{str(phase_df[phase_df['PRT 401']==10]['HYS 501'].max())} %"
                        for row in table.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                elif phase == 'Reverse BPR':
                    phase_df["PRT 401"] = phase_df["PRT 401"].apply(lambda x: round(float(x)))
                    phase_df['Current SP'] = phase_df['Current SP'].astype(int)
                    #print(phase_df.to_csv('bpr217valuecheck.csv'))
                    fil = phase_df[phase_df['PRT 401']==10]
                    if not fil.empty:
                        num_cols = phase_df.select_dtypes(include=['number'])
                        last_row = num_cols[num_cols['PRT 401']==10].mean().round(3)
                        max_current = table.table.cell(1,5)
                        max_current.text = f"{str(phase_df['Current SP'].max())} A"
                        cell_columns = last_row.filter(like='Cell')
                        cell_col = cell_columns.astype(float)
                        max_cell_value = cell_col.max()
                        max_cell_vol = table.table.cell(2,5)
                        max_cell_vol.text = f"{str(max_cell_value)} V"
                        max_cell_col = cell_col.idxmax()
                        max_cell_column = table.table.cell(3,5)
                        max_cell_column.text = str(max_cell_col)
                        min_cell_value = cell_col.min()
                        min_cell_vol = table.table.cell(4,5)
                        min_cell_vol.text = f"{str(min_cell_value)} V"
                        min_cell_col = cell_col.idxmin()
                        max_cell_column = table.table.cell(5,5)
                        max_cell_column.text = str(min_cell_col)
                        avg_vol = table.table.cell(6,5)
                        avg_vol.text = f"{str(phase_df[phase_df['PRT 401']==10]['avg_CV'].mean().round(3))} V"
                        water_temp = table.table.cell(7,5)
                        water_temp.text = f"{str(phase_df[phase_df['PRT 401']==10]['TTC 102'].mean().round(1))} °C"
                        inlet_h2_pre = table.table.cell(8,5)
                        inlet_h2_pre.text = f"{str(phase_df[phase_df['PRT 401']==10]['PRT 104'].astype(float).mean().round(2))} bar"
                        cos = table.table.cell(9,5)
                        cos.text = f"{str(phase_df[phase_df['PRT 401']==10]['COS 101'].mean().round(1))} µS/cm"
                        inherent_pre = table.table.cell(10,5)
                        inherent_pre.text= f"{str(phase_df[phase_df['PRT 401']==10]['PRT 401'].mean().round(2))} bar"
                        hys1 = table.table.cell(11,5)
                        hys1.text_frame.text = f"{str(phase_df[phase_df['PRT 401']==10]['HYS 102'].max())} %"
                        hys2 = table.table.cell(12,5)
                        hys2.text_frame.text = f"{str(phase_df[phase_df['PRT 401']==10]['HYS 102'][-10:].mean())} %"
                        hys3 = table.table.cell(13,5)
                        hys3.text_frame.text = f"{str(phase_df[phase_df['PRT 401']==10]['HYS 501'].max())} %"
                        for row in table.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                        
                    
            
    slide = prs.slides[1] 
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(9.0) 
    # place the image
    pic = slide.shapes.add_picture('{}/back_pressure_test.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/back_pressure_test.html'.format(image_folder_name)

    slide = prs.slides[2]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(9.0) 
    # place the image
    pic = slide.shapes.add_picture('{}/polarization_chart.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/polarization_chart.html'.format(image_folder_name)


    slide = prs.slides[3]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(9.0) 
    # place the image
    pic = slide.shapes.add_picture('{}/steady_state_(840a)_voltage_distribution.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/steady_state_(840a)_voltage_distribution.html'.format(image_folder_name)
      
    slide = prs.slides[4]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(9.0) 
    # place the image
    pic = slide.shapes.add_picture('{}/10bar_voltage_distribution.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10bar_voltage_distribution.html'.format(image_folder_name)

    slide = prs.slides[5]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(9.0) 
    # place the image
    pic = slide.shapes.add_picture('{}/13bar_voltage_distribution.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/13bar_voltage_distribution.html'.format(image_folder_name)
      

    slide = prs.slides[6]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(9.0) 
    # place the image
    pic = slide.shapes.add_picture('{}/steady_state_voltage_distribution_by_pressure_(all_cells).png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/steady_state_voltage_distribution_by_pressure_(all_cells).html'.format(image_folder_name)
    
    slide = prs.slides[7]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(9.0) 
    # place the image
    pic = slide.shapes.add_picture('{}/ocv_decay_chart.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/ocv_decay_chart.png'.format(image_folder_name)
    
    slide = prs.slides[8]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the steady state image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle1.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle1.html'.format(image_folder_name)
    slide = prs.slides[9]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the operating condition image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle2.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle2.html'.format(image_folder_name)
    slide = prs.slides[10]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the ASR image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle3.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle3.html'.format(image_folder_name)
    slide = prs.slides[11]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the hfr and flow image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle4.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle4.html'.format(image_folder_name)
    
    slide = prs.slides[12]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the steady state image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle5.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle5.html'.format(image_folder_name)
    slide = prs.slides[13]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the operating condition image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle6.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle6.html'.format(image_folder_name)
    slide = prs.slides[14]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the ASR image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle7.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle7.html'.format(image_folder_name)
    slide = prs.slides[15]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the hfr and flow image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle8.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle8.html'.format(image_folder_name)
    
    slide = prs.slides[16]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the steady state image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle9.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle9.html'.format(image_folder_name)
    slide = prs.slides[17]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the operating condition image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle10.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle10.html'.format(image_folder_name)
    slide = prs.slides[18]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the ASR image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle11.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle11.html'.format(image_folder_name)
    slide = prs.slides[19]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the hfr and flow image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle12.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle12.html'.format(image_folder_name)
    
    slide = prs.slides[20]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the steady state image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle13.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle13.html'.format(image_folder_name)
    slide = prs.slides[21]
    #set the height for the images
    row = Inches(0.2)
    column = Inches(0.2)
    height = Inches(7.0)
    width = Inches(10.0) 
    # place the operating condition image
    pic = slide.shapes.add_picture('{}/10_bar_voltage_distribution_for_bundle14.png'.format(image_folder_name),row, column, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/10_bar_voltage_distribution_for_bundle14.html'.format(image_folder_name)

    # change to the output directory
    prs.save('{}/{}_{}.pptx'.format(output_dir, data['Stack ID'][0], 'PST_Report'))
    

def main_fun(file_path, output_dir):

    global image_folder_name
    image_folder_name = '{}/images'.format(output_dir)
    if not os.path.exists(image_folder_name):
        os.mkdir(image_folder_name)
    #Read the file
    print('Reading Dataframe')
    data = read_files(file_path)
    #Define number of polarization test performed with column 'poltest'
    new_df = pol_test(data)
    #print(new_df)
    final_df = raw_data(new_df)
    #print(final_df)
    filtered_df = processing_final_df(final_df)
    #print(filtered_df)
    #Finalised DataFrame to process requried plots
    print('Finalized DataFrame')
    filtered_df = processing_final_df(final_df)
    #Plot polarization chart
    print('Plotting all the data')
    ovr_all = plotting_overall_data(filtered_df)
    print(ovr_all)
    #Plot BPR chart
    print('Plotting BPR Chart')
    bpr = plotting_BPR_plot(filtered_df)
    print(bpr)
    #Plot polarization chart
    print('Plotting polarization chart')
    pol_chart = plotting_polarization_data(filtered_df)
    print(pol_chart)
    #Plot steady state 840A voltage distribution plot (by bundles)
    print('Plotting Steady State (840A) Voltage Distribution Chart')
    ss_vol_dis = plotting_steady_state_vol_dis(filtered_df)
    print(ss_vol_dis)
    #Plot voltage distribution at different pressures (all cells)
    print('Plotting Steady State Voltage Distribution by Preessure Chart')
    ss_vol_dis_by_p = steady_state_vol_dis_by_pressure(filtered_df)
    print(ss_vol_dis_by_p)
    #Plot 10bar voltage distribution plot
    print('Plotting 10bar Voltage Distribution by Bundles chart')
    bp_10bar_vol_dis = back_pressure_by_bundle10(filtered_df)
    print(bp_10bar_vol_dis)
    bp_13bar_vol_dis = back_pressure_by_bundle13(filtered_df)
    print('plotting ocv decay chart')
    ocv_chart =  plot_OCV_decay(filtered_df)
    #print(ocv_chart)
    bundle_1 = plotting_bundle1_vol_dis(filtered_df)
    print(bundle_1)
    bundle_2 = plotting_bundle2_vol_dis(filtered_df)
    print(bundle_2)
    bundle_3 = plotting_bundle3_vol_dis(filtered_df)
    print(bundle_3)
    bundle_4 = plotting_bundle4_vol_dis(filtered_df)
    print(bundle_4)
    bundle_5 = plotting_bundle5_vol_dis(filtered_df)
    print(bundle_5)
    bundle_6 = plotting_bundle6_vol_dis(filtered_df)
    print(bundle_6)
    bundle_7 = plotting_bundle7_vol_dis(filtered_df)
    print(bundle_7)
    bundle_8 = plotting_bundle8_vol_dis(filtered_df)
    print(bundle_8)
    bundle_9 = plotting_bundle9_vol_dis(filtered_df)
    print(bundle_9)
    bundle_10 = plotting_bundle10_vol_dis(filtered_df)
    print(bundle_10)
    bundle_11 = plotting_bundle11_vol_dis(filtered_df)
    print(bundle_11)
    bundle_12 = plotting_bundle12_vol_dis(filtered_df)
    print(bundle_12)
    bundle_13 = plotting_bundle13_vol_dis(filtered_df)
    print(bundle_13)
    bundle_14 = plotting_bundle14_vol_dis(filtered_df)
    print(bundle_14)
    #Outputting the pptx
    create_pptx(file_path, image_folder_name, data, filtered_df, final_df, output_dir)
    print("Outputting pptx")   
    print('Program Complete')  



#%% This cell defines the function for running the gui
def rungui():
    
    # Create an instance of tkinter frame
    win = Tk()
    win.title('PST Data Visualization')
    
    # create the frame layout for the main window
    frame = tk.Frame(master=win, width=500, height=175)
    frame.pack()
    
    # add a label for the automatic data
    label1 = tk.Label(master=frame, text="Raw data file:")
    label1.place(x=10, y=25)
    
    # add an entry box for the automatic data file
    entry1 = tk.Entry(win, bg="white", width = 45)
    entry1.place(x=125, y=25)
    
    # this function opens a file browser for excel files and outputs the selected box to the automatic data entry field
    def browse1():
        try:
            entry1.delete(0, tk.END)
            files = filedialog.askopenfilenames(filetypes=[('Excel and CSV Files', '*.xlsx *.csv')])
            for i in files[:-1]:
                entry1.insert(tk.END, '{}, '.format(i))
            entry1.insert(tk.END, files[-1])
        except PermissionError:
            messagebox.showerror('Permission Error', 'Error: Please close the file on your computer and try again.')
    
    # Create a browse button that links to the browse1 function
    button1 = ttk.Button(win, text="Browse", command=browse1)
    button1.place(x=400, y=21)
    
    #create a label for the output folder
    label3 = tk.Label(master=frame, text="Output folder:")
    label3.place(x=37, y=77)
    
    # creat an entry box for the output folder    
    entry2 = tk.Entry(win, bg="white", width = 45)
    entry2.place(x=125, y=77)
    
    # this function opens a file browser for excel files and outputs the 
    # selected box to the manual data entry field
    def browse2():
        try:
            folder = filedialog.askdirectory()
            #foldername = os.path.abspath(folder.name)
            entry2.insert(tk.END, folder)
        except PermissionError:
            messagebox.showerror('Permission Error', 'Error: Please close the folder on your computer and try again.')
    
    # Create a browse button that links to the browse3 function
    button2 = ttk.Button(win, text="Browse", command=browse2)
    button2.place(x=400, y=74)
    
    # this defines the main run function that we wish to do
    def run_main():
        file_path = entry1.get().split(', ')
        output_dir = entry2.get()
        main_fun(file_path, output_dir)
        win.destroy()
        
    # this is the overall function called that does the main function and the progress bar
    def run_function(name, func):
        # disable the run button
        button4['state'] = 'disabled'
        # define and start the progress bar
        progress_bar = ttk.Progressbar(win, orient = HORIZONTAL, length = 400, mode='indeterminate')
        progress_bar.place(x=50, y = 115)
        progress_bar.start(interval = 15)
        # run the main function
        func()
        # stop the progress bar
        progress_bar.stop()
        # print program complete
        print('Program complete')
        #destory the progress bar
        progress_bar.destroy()
        # show the finished label
        label4 = tk.Label(master=frame, text="Finished!")
        label4.place(x=240, y=115)
        # renable the run button
        button4['state'] = 'enabled'
        
    # this is the function that threads together the above two functions
    def run_thread(name, func):
        Thread(target = run_function, args=(name,func)).start()
    
    # this is the function that is called when the run button is clicked
    def run_clicked():
        run_thread('main_fun', run_main)
    
    # Create a Button that runs the main script
    button4 = ttk.Button(win, text="Run", command=run_clicked)
    button4.place(x=225, y=145)
    
    # create a help menu
                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                            
    win.mainloop()

#%% Run the GUI
rungui()